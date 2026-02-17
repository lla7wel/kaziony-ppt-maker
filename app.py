import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

W, H = Inches(13.333), Inches(7.5)

COL = {
    "bg": RGBColor(10, 10, 14),
    "panel": RGBColor(18, 18, 24),
    "panel2": RGBColor(24, 24, 32),
    "gold": RGBColor(212, 175, 55),
    "white": RGBColor(245, 245, 247),
    "muted": RGBColor(170, 170, 178),
    "line": RGBColor(60, 60, 72),
    "ok": RGBColor(68, 214, 144),
    "warn": RGBColor(255, 204, 0),
}

FONT_EN = "Segoe UI"
FONT_AR = "Segoe UI Arabic"

def add_bg(slide):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid()
    r.fill.fore_color.rgb = COL["bg"]
    r.line.fill.background()

def add_header(slide, title_en, title_ar, app_name="Kaziony"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COL["panel"]
    bar.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.82), W, Inches(0.03))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COL["gold"]
    accent.line.fill.background()

    tb_en = slide.shapes.add_textbox(Inches(0.65), Inches(0.18), Inches(6.2), Inches(0.55))
    tf = tb_en.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_en
    p.font.name = FONT_EN
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COL["white"]
    p.alignment = PP_ALIGN.LEFT

    tb_ar = slide.shapes.add_textbox(Inches(6.9), Inches(0.18), Inches(5.8), Inches(0.55))
    tf = tb_ar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_ar
    p.font.name = FONT_AR
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COL["white"]
    p.alignment = PP_ALIGN.RIGHT

    brand = slide.shapes.add_textbox(Inches(11.7), Inches(0.22), Inches(1.55), Inches(0.4))
    tf = brand.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = app_name
    p.font.name = FONT_EN
    p.font.size = Pt(14)
    p.font.color.rgb = COL["muted"]
    p.alignment = PP_ALIGN.RIGHT

def add_footer(slide, idx, total):
    tb = slide.shapes.add_textbox(Inches(0.65), Inches(7.18), Inches(12.0), Inches(0.3))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{idx}/{total}"
    p.font.name = FONT_EN
    p.font.size = Pt(12)
    p.font.color.rgb = COL["muted"]
    p.alignment = PP_ALIGN.RIGHT

def panel(slide, x, y, w, h, accent=None, lw=1):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = COL["panel"]
    s.line.color.rgb = accent if accent else COL["line"]
    s.line.width = Pt(lw)
    return s

def chip(slide, x, y, text, font=FONT_EN, fg=COL["bg"], bgc=COL["gold"], w=Inches(2.2), h=Inches(0.42), size=14):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = bgc
    s.line.fill.background()
    tf = s.text_frame; tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = fg
    p.alignment = PP_ALIGN.CENTER
    return s

def add_bilingual_cards(slide, items, x, y, w, h, cols=2):
    # items: (en_title, ar_title, en_desc, ar_desc, accent_rgb)
    rows = math.ceil(len(items) / cols)
    gapx, gapy = Inches(0.35), Inches(0.28)
    card_w = (w - gapx*(cols-1)) / cols
    card_h = (h - gapy*(rows-1)) / rows
    for i, it in enumerate(items):
        r = i // cols
        c = i % cols
        cx = x + c * (card_w + gapx)
        cy = y + r * (card_h + gapy)
        accent = it[4] if len(it) > 4 and it[4] is not None else COL["gold"]
        panel(slide, cx, cy, card_w, card_h, accent=accent, lw=2)
        tb = slide.shapes.add_textbox(cx+Inches(0.28), cy+Inches(0.18), card_w-Inches(0.56), card_h-Inches(0.32))
        tf = tb.text_frame; tf.clear(); tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = it[0]
        p.font.name = FONT_EN; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = it[1]
        p2.font.name = FONT_AR; p2.font.size = Pt(16); p2.font.bold = True; p2.font.color.rgb = COL["white"]
        p2.alignment = PP_ALIGN.RIGHT

        p3 = tf.add_paragraph()
        p3.text = it[2]
        p3.font.name = FONT_EN; p3.font.size = Pt(13); p3.font.color.rgb = COL["muted"]

        p4 = tf.add_paragraph()
        p4.text = it[3]
        p4.font.name = FONT_AR; p4.font.size = Pt(13); p4.font.color.rgb = COL["muted"]
        p4.alignment = PP_ALIGN.RIGHT

def add_three_step(slide):
    items = [
        ("Order", "الطلب", "Customer selects merchant → items → checkout", "العميل يختار المتجر → يضيف المنتجات → تأكيد", COL["gold"]),
        ("Credit", "كريديت", "Driver accepts using 1 credit (deducted instantly)", "السائق يقبل بكريديت واحد (خصم فوري)", COL["ok"]),
        ("Cash", "كاش", "Deliver → collect cash (items + delivery + tip)", "تسليم → تحصيل الكاش (قيمة الطلب + التوصيل + التيب)", COL["warn"]),
    ]
    xs = [Inches(0.9), Inches(4.8), Inches(8.7)]
    for i, it in enumerate(items):
        panel(slide, xs[i], Inches(2.0), Inches(3.8), Inches(2.6), accent=it[4], lw=2)
        tb = slide.shapes.add_textbox(xs[i]+Inches(0.3), Inches(2.2), Inches(3.2), Inches(2.2))
        tf = tb.text_frame; tf.clear(); tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{it[0]} | {it[1]}"
        p.font.name = FONT_EN; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = it[4]
        p2 = tf.add_paragraph()
        p2.text = it[2]
        p2.font.name = FONT_EN; p2.font.size = Pt(14); p2.font.color.rgb = COL["white"]
        p3 = tf.add_paragraph()
        p3.text = it[3]
        p3.font.name = FONT_AR; p3.font.size = Pt(14); p3.font.color.rgb = COL["white"]; p3.alignment = PP_ALIGN.RIGHT
        if i < 2:
            arr = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, xs[i]+Inches(3.9), Inches(3.0), Inches(0.8), Inches(0.75))
            arr.fill.solid(); arr.fill.fore_color.rgb = COL["gold"]; arr.line.fill.background()

def add_flow(slide):
    steps_en = ["Placed", "Merchant accepts", "Preparing", "Ready", "Driver accepts\n(1 credit)", "Delivered\n+ cash"]
    steps_ar = ["تأكيد الطلب", "قبول المتجر", "تجهيز", "جاهز", "قبول السائق\n(كريديت 1)", "تسليم\n+ كاش"]
    x0, y, w, h, gap = Inches(0.8), Inches(2.2), Inches(1.95), Inches(1.3), Inches(0.23)
    for i in range(len(steps_en)):
        x = x0 + i*(w+gap)
        panel(slide, x, y, w, h)
        chip(slide, x+Inches(0.25), y-Inches(0.38), f"{i+1}", FONT_EN, COL["bg"], COL["gold"], w=Inches(0.5), h=Inches(0.35), size=12)
        tb = slide.shapes.add_textbox(x+Inches(0.15), y+Inches(0.20), w-Inches(0.3), h-Inches(0.3))
        tf = tb.text_frame; tf.clear(); tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = steps_en[i]
        p.font.name = FONT_EN; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COL["gold"]
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = steps_ar[i]
        p2.font.name = FONT_AR; p2.font.size = Pt(14); p2.font.color.rgb = COL["white"]
        p2.alignment = PP_ALIGN.CENTER
        if i < len(steps_en)-1:
            ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+w, y+h/2, x+w+gap, y+h/2)
            ln.line.color.rgb = COL["gold"]; ln.line.width = Pt(2)

def add_phone(slide, x, y, title_en, title_ar, cta_text, kind="generic"):
    phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.7), Inches(5.0))
    phone.fill.solid(); phone.fill.fore_color.rgb = COL["panel"]
    phone.line.color.rgb = COL["gold"]; phone.line.width = Pt(2)

    screen = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x+Inches(0.18), y+Inches(0.50), Inches(2.34), Inches(4.3))
    screen.fill.solid(); screen.fill.fore_color.rgb = COL["panel2"]; screen.line.fill.background()

    tb = slide.shapes.add_textbox(x+Inches(0.25), y+Inches(0.62), Inches(2.2), Inches(0.65))
    tf = tb.text_frame; tf.clear()
    p = tf.paragraphs[0]
    p.text = title_en
    p.font.name = FONT_EN; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = COL["gold"]
    p2 = tf.add_paragraph()
    p2.text = title_ar
    p2.font.name = FONT_AR; p2.font.size = Pt(13); p2.font.color.rgb = COL["white"]; p2.alignment = PP_ALIGN.RIGHT

    if kind == "tracking":
        m = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Inches(0.28), y+Inches(1.38), Inches(2.14), Inches(2.1))
        m.fill.solid(); m.fill.fore_color.rgb = COL["panel"]; m.line.color.rgb = COL["line"]
        l = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+Inches(0.45), y+Inches(2.0), x+Inches(2.2), y+Inches(2.8))
        l.line.color.rgb = COL["gold"]; l.line.width = Pt(2)
        chip(slide, x+Inches(0.35), y+Inches(3.6), "En route", FONT_EN, COL["bg"], COL["gold"], w=Inches(1.0), h=Inches(0.33), size=11)
        chip(slide, x+Inches(1.45), y+Inches(3.6), "في الطريق", FONT_AR, COL["bg"], COL["gold"], w=Inches(1.0), h=Inches(0.33), size=11)
    elif kind == "checkout":
        chip(slide, x+Inches(0.35), y+Inches(1.45), "Fee shown upfront", FONT_EN, COL["bg"], COL["gold"], w=Inches(1.9), h=Inches(0.33), size=11)
        chip(slide, x+Inches(0.35), y+Inches(1.85), "الرسوم تظهر قبل التأكيد", FONT_AR, COL["bg"], COL["gold"], w=Inches(1.9), h=Inches(0.33), size=11)
        for k in range(3):
            b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Inches(0.28), y+Inches(2.35 + k*0.72), Inches(2.14), Inches(0.55))
            b.fill.solid(); b.fill.fore_color.rgb = COL["panel"]; b.line.color.rgb = COL["line"]
    elif kind == "offer":
        chip(slide, x+Inches(0.35), y+Inches(1.45), "New offer", FONT_EN, COL["bg"], COL["ok"], w=Inches(1.0), h=Inches(0.33), size=11)
        chip(slide, x+Inches(1.42), y+Inches(1.45), "عرض جديد", FONT_AR, COL["bg"], COL["ok"], w=Inches(1.0), h=Inches(0.33), size=11)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Inches(0.28), y+Inches(1.95), Inches(2.14), Inches(1.6))
        card.fill.solid(); card.fill.fore_color.rgb = COL["panel"]; card.line.color.rgb = COL["line"]
        chip(slide, x+Inches(0.38), y+Inches(3.65), "1 credit to accept", FONT_EN, COL["bg"], COL["gold"], w=Inches(2.0), h=Inches(0.33), size=11)
    elif kind == "pickup":
        for k in range(4):
            b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Inches(0.28), y+Inches(1.45 + k*0.68), Inches(2.14), Inches(0.52))
            b.fill.solid(); b.fill.fore_color.rgb = COL["panel"]; b.line.color.rgb = COL["line"]
        chip(slide, x+Inches(0.38), y+Inches(4.25), "Receipt photo", FONT_EN, COL["bg"], COL["gold"], w=Inches(1.0), h=Inches(0.33), size=11)
        chip(slide, x+Inches(1.52), y+Inches(4.25), "إيصال", FONT_AR, COL["bg"], COL["gold"], w=Inches(0.9), h=Inches(0.33), size=11)
    else:
        for k in range(4):
            b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Inches(0.28), y+Inches(1.45 + k*0.78), Inches(2.14), Inches(0.55))
            b.fill.solid(); b.fill.fore_color.rgb = COL["panel"]; b.line.color.rgb = COL["line"]

    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Inches(0.60), y+Inches(4.5), Inches(1.5), Inches(0.42))
    btn.fill.solid(); btn.fill.fore_color.rgb = COL["gold"]; btn.line.fill.background()
    tf = btn.text_frame; tf.clear()
    p = tf.paragraphs[0]
    p.text = cta_text
    p.font.name = FONT_EN; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = COL["bg"]
    p.alignment = PP_ALIGN.CENTER

def add_money_flow(slide):
    y = Inches(1.95)
    nodes = [("Merchant", "المتجر", Inches(0.95)), ("Driver", "السائق", Inches(5.55)), ("Customer", "العميل", Inches(10.15))]
    for en, ar, x in nodes:
        panel(slide, x, y, Inches(2.25), Inches(1.25))
        tb = slide.shapes.add_textbox(x+Inches(0.25), y+Inches(0.22), Inches(1.75), Inches(0.9))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = en; p.font.name = FONT_EN; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COL["gold"]; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = ar; p2.font.name = FONT_AR; p2.font.size = Pt(16); p2.font.color.rgb = COL["white"]; p2.alignment = PP_ALIGN.CENTER

    def arrow(x1, x2, label_en, label_ar):
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y+Inches(0.62), x2, y+Inches(0.62))
        ln.line.color.rgb = COL["gold"]; ln.line.width = Pt(3)
        tb = slide.shapes.add_textbox((x1+x2)/2 - Inches(1.35), y+Inches(1.35), Inches(2.7), Inches(0.8))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = label_en; p.font.name = FONT_EN; p.font.size = Pt(12); p.font.color.rgb = COL["white"]; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = label_ar; p2.font.name = FONT_AR; p2.font.size = Pt(12); p2.font.color.rgb = COL["white"]; p2.alignment = PP_ALIGN.CENTER

    arrow(Inches(3.2), Inches(5.55), "Driver pays items (pickup)", "السائق يدفع قيمة الطلب (استلام)")
    arrow(Inches(7.85), Inches(10.15), "Customer pays: items + delivery + tip", "العميل يدفع: الطلب + التوصيل + التيب")

    cards = [
        ("Driver earnings", "دخل السائق", "Delivery fee + tips", "رسوم التوصيل + التيب", COL["gold"]),
        ("Items cash", "قيمة الطلب", "Reimbursement only (passes through)", "استرجاع فقط (تمرير)", COL["gold"]),
    ]
    add_bilingual_cards(slide, cards, Inches(0.95), Inches(4.6), Inches(12.0), Inches(2.1), cols=2)

def add_dispatch_visual(slide):
    base_y = Inches(2.05)
    panel(slide, Inches(0.9), base_y, Inches(3.7), Inches(1.2))
    tb = slide.shapes.add_textbox(Inches(1.15), base_y+Inches(0.18), Inches(3.2), Inches(0.9))
    tf = tb.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.text = "Auto-offer to nearby drivers"; p.font.name = FONT_EN; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COL["gold"]
    p2 = tf.add_paragraph(); p2.text = "عرض تلقائي للسائقين القريبين"; p2.font.name = FONT_AR; p2.font.size = Pt(16); p2.font.color.rgb = COL["white"]; p2.alignment = PP_ALIGN.RIGHT

    panel(slide, Inches(5.1), base_y-Inches(0.35), Inches(3.6), Inches(1.0))
    tf = slide.shapes.add_textbox(Inches(5.35), base_y-Inches(0.18), Inches(3.1), Inches(0.7)).text_frame
    tf.clear()
    p = tf.paragraphs[0]; p.text = "Driver accepts"; p.font.name = FONT_EN; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COL["ok"]; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "سائق يقبل"; p2.font.name = FONT_AR; p2.font.size = Pt(15); p2.font.color.rgb = COL["white"]; p2.alignment = PP_ALIGN.CENTER

    panel(slide, Inches(5.1), base_y+Inches(0.95), Inches(3.6), Inches(1.0))
    tf = slide.shapes.add_textbox(Inches(5.35), base_y+Inches(1.12), Inches(3.1), Inches(0.7)).text_frame
    tf.clear()
    p = tf.paragraphs[0]; p.text = "No driver found"; p.font.name = FONT_EN; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COL["warn"]; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "لا يوجد سائق"; p2.font.name = FONT_AR; p2.font.size = Pt(15); p2.font.color.rgb = COL["white"]; p2.alignment = PP_ALIGN.CENTER

    panel(slide, Inches(9.25), base_y+Inches(0.95), Inches(3.15), Inches(1.0))
    tf = slide.shapes.add_textbox(Inches(9.45), base_y+Inches(1.12), Inches(2.75), Inches(0.7)).text_frame
    tf.clear()
    p = tf.paragraphs[0]; p.text = "Customer option: Wait longer"; p.font.name = FONT_EN; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = COL["gold"]; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "خيار العميل: انتظار أكثر"; p2.font.name = FONT_AR; p2.font.size = Pt(13); p2.font.color.rgb = COL["white"]; p2.alignment = PP_ALIGN.CENTER

    def conn(x1, y1, x2, y2):
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        ln.line.color.rgb = COL["gold"]; ln.line.width = Pt(2)

    conn(Inches(4.6), base_y+Inches(0.6), Inches(5.1), base_y+Inches(0.15))
    conn(Inches(4.6), base_y+Inches(0.6), Inches(5.1), base_y+Inches(1.45))
    conn(Inches(8.7), base_y+Inches(1.45), Inches(9.25), base_y+Inches(1.45))

def add_admin_dashboard(slide):
    kpis = [
        ("Orders/day", "طلبات/يوم"),
        ("Completion", "نسبة الإكمال"),
        ("Avg ETA", "متوسط الوقت"),
        ("Active drivers", "سائقون متصلون"),
        ("Credits sold", "كريديت مباعة"),
        ("Credits used", "كريديت مستخدمة"),
    ]
    x0, y0 = Inches(7.1), Inches(1.55)
    w, h = Inches(2.05), Inches(0.95)
    gapx, gapy = Inches(0.25), Inches(0.22)
    for i, (en, ar) in enumerate(kpis):
        cx = x0 + (i % 2) * (w + gapx)
        cy = y0 + (i // 2) * (h + gapy)
        panel(slide, cx, cy, w, h)
        tb = slide.shapes.add_textbox(cx+Inches(0.18), cy+Inches(0.14), w-Inches(0.36), h-Inches(0.2))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = en; p.font.name = FONT_EN; p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = COL["gold"]
        p2 = tf.add_paragraph(); p2.text = ar; p2.font.name = FONT_AR; p2.font.size = Pt(11); p2.font.color.rgb = COL["white"]; p2.alignment = PP_ALIGN.RIGHT
        p3 = tf.add_paragraph(); p3.text = "—"; p3.font.name = FONT_EN; p3.font.size = Pt(18); p3.font.bold = True; p3.font.color.rgb = COL["white"]

    chart_x, chart_y = Inches(7.1), Inches(5.4)
    chart_w, chart_h = Inches(5.25), Inches(1.95)
    panel(slide, chart_x, chart_y, chart_w, chart_h)
    heights = [0.55, 1.05, 0.8, 1.45, 0.9, 1.65, 1.1]
    for i, hh in enumerate(heights):
        b = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            chart_x+Inches(0.45+i*0.7),
            chart_y+chart_h-Inches(0.25)-Inches(hh),
            Inches(0.35),
            Inches(hh),
        )
        b.fill.solid(); b.fill.fore_color.rgb = COL["gold"]; b.line.fill.background()
    tb = slide.shapes.add_textbox(chart_x+Inches(0.25), chart_y+Inches(0.15), Inches(3.0), Inches(0.35))
    tf = tb.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.text = "Example trend (placeholder)"; p.font.name = FONT_EN; p.font.size = Pt(10); p.font.color.rgb = COL["muted"]

def build_fixed_deck(app_name="Kaziony", city_en="Tripoli", city_ar="طرابلس"):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    total = 12
    idx = 1

    # 1 Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    title = s.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.6), Inches(2.2))
    tf = title.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"{app_name} | كازيوني"; p.font.name = FONT_EN; p.font.size = Pt(56); p.font.bold = True; p.font.color.rgb = COL["gold"]
    p2 = tf.add_paragraph(); p2.text = "Investor overview • Cash-first delivery + driver-credit acceptance"; p2.font.name = FONT_EN; p2.font.size = Pt(20); p2.font.color.rgb = COL["white"]
    p3 = tf.add_paragraph(); p3.text = "عرض للمستثمر • توصيل كاش + قبول الطلب بكريديت للسائق"; p3.font.name = FONT_AR; p3.font.size = Pt(20); p3.font.color.rgb = COL["white"]; p3.alignment = PP_ALIGN.RIGHT
    chip(s, Inches(0.9), Inches(6.2), f"Launch: {city_en} / {city_ar}", FONT_EN, COL["bg"], COL["gold"], w=Inches(3.8), h=Inches(0.45))
    add_footer(s, idx, total); idx += 1

    # 2 What it is
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_header(s, "What is Kaziony", "ما هو كازيوني", app_name)
    chip(s, Inches(0.9), Inches(1.15), "Cash on Delivery", FONT_EN, COL["bg"], COL["gold"], w=Inches(2.6))
    chip(s, Inches(3.65), Inches(1.15), "Fees shown upfront", FONT_EN, COL["bg"], COL["gold"], w=Inches(2.9))
    chip(s, Inches(6.7), Inches(1.15), "Live tracking", FONT_EN, COL["bg"], COL["gold"], w=Inches(2.0))
    chip(s, Inches(8.85), Inches(1.15), "Hybrid dispatch", FONT_EN, COL["bg"], COL["gold"], w=Inches(2.3))
    items = [
        ("Restaurants", "مطاعم", "Order food from nearby merchants", "طلب الطعام من المطاعم القريبة", COL["gold"]),
        ("Groceries", "بقالة", "Daily items and essentials", "احتياجات يومية وأساسية", COL["gold"]),
        ("Pharmacy", "صيدلية", "Fast pharmacy deliveries", "توصيل الصيدلية بسرعة", COL["gold"]),
        ("Scheduled", "مجدول", "Choose a delivery time window", "حدد نافذة زمنية للتوصيل", COL["gold"]),
    ]
    add_bilingual_cards(s, items, Inches(0.9), Inches(1.75), Inches(12.0), Inches(5.2), cols=2)
    add_footer(s, idx, total); idx += 1

    # 3 Steps
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_header(s, "How it works (3 steps)", "كيف يعمل (٣ خطوات)", app_name)
    add_three_step(s)
    add_footer(s, idx, total); idx += 1

    # 4 Flow
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_header(s, "Order flow (end-to-end)", "مسار الطلب (من البداية للنهاية)", app_name)
    add_flow(s)
    add_footer(s, idx, total); idx += 1

    # 5 Customer journey
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_header(s, "Customer journey (visual)", "رحلة العميل (مرئي)", app_name)
    add_phone(s, Inches(0.95), Inches(1.55), "Browse", "تصفح", "Add", kind="generic")
    add_phone(s, Inches(4.15), Inches(1.55), "Cart", "السلة", "Checkout", kind="generic")
    add_phone(s, Inches(7.35), Inches(1.55), "Checkout", "الدفع", "Confirm", kind="checkout")
    add_phone(s, Inches(10.55), Inches(1.55), "Tracking", "التتبع", "Chat", kind="tracking")
    add_footer(s, idx, total); idx += 1

    # 6 Driver journey
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_header(s, "Driver journey (visual)", "رحلة السائق (مرئي)", app_name)
    add_phone(s, Inches(0.95), Inches(1.55), "Offers", "عروض", "Open", kind="offer")
    add_phone(s, Inches(4.15), Inches(1.55), "Accept", "قبول", "Accept", kind="offer")
    add_phone(s, Inches(7.35), Inches(1.55), "Pickup", "استلام", "Paid", kind="pickup")
    add_phone(s, Inches(10.55), Inches(1.55), "Deliver", "تسليم", "Done", kind="tracking")
    add_footer(s, idx, total); idx += 1

    # 7 Credits
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_header(s, "Credits (core mechanic)", "الكريديت (الآلية الأساسية)", app_name)
    rules = [
        ("1 credit per acceptance", "كريديت واحد لكل قبول", "Deducted instantly on 'Accept'", "خصم فوري عند القبول", COL["gold"]),
        ("Refund on cancellation", "استرجاع عند الإلغاء", "Credit returns when order cancels", "يرجع الكريديت عند الإلغاء", COL["ok"]),
        ("Top-up method", "طريقة الشحن", "In-app wallet top-up (planned)", "شحن محفظة داخل التطبيق (مخطط)", COL["warn"]),
        ("Why credits", "لماذا كريديت", "Reduces spam accepts and improves reliability", "يقلل القبول العشوائي ويزيد الموثوقية", COL["gold"]),
    ]
    add_bilingual_cards(s, rules, Inches(0.9), Inches(1.55), Inches(12.0), Inches(5.4), cols=2)
    add_footer(s, idx, total); idx += 1

    # 8 Top-up steps
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_header(s, "Credits: top-up steps", "الكريديت: خطوات الشحن", app_name)
    steps = [
        ("Top up wallet", "شحن المحفظة", "Inside the app (planned)", "داخل التطبيق (مخطط)"),
        ("Balance", "الرصيد", "Credits available", "الكريديت جاهزة"),
        ("Accept order", "قبول الطلب", "1 credit deducted", "خصم كريديت 1"),
        ("Cancel", "إلغاء", "Credit refunded", "استرجاع الكريديت"),
    ]
    x0, y, w, h, gap = Inches(0.95), Inches(2.45), Inches(2.8), Inches(1.75), Inches(0.35)
    for i, (a, b, c, d) in enumerate(steps):
        x = x0 + i*(w+gap)
        panel(s, x, y, w, h, accent=COL["gold"], lw=2)
        chip(s, x+Inches(0.15), y-Inches(0.45), f"{i+1}", FONT_EN, COL["bg"], COL["gold"], w=Inches(0.5), h=Inches(0.38), size=12)
        tb = s.shapes.add_textbox(x+Inches(0.22), y+Inches(0.18), w-Inches(0.44), h-Inches(0.25))
        tf = tb.text_frame; tf.clear(); tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = a; p.font.name = FONT_EN; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COL["gold"]; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = b; p2.font.name = FONT_AR; p2.font.size = Pt(15); p2.font.color.rgb = COL["white"]; p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph(); p3.text = c; p3.font.name = FONT_EN; p3.font.size = Pt(12); p3.font.color.rgb = COL["muted"]; p3.alignment = PP_ALIGN.CENTER
        p4 = tf.add_paragraph(); p4.text = d; p4.font.name = FONT_AR; p4.font.size = Pt(12); p4.font.color.rgb = COL["muted"]; p4.alignment = PP_ALIGN.CENTER
        if i < 3:
            ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x+w, y+h/2, x+w+gap, y+h/2)
            ln.line.color.rgb = COL["gold"]; ln.line.width = Pt(2)
    add_footer(s, idx, total); idx += 1

    # 9 Money flow
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_header(s, "Money flow (cash)", "مسار الأموال (كاش)", app_name)
    add_money_flow(s)
    add_footer(s, idx, total); idx += 1

    # 10 Pricing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_header(s, "Pricing (distance-based)", "التسعير (حسب المسافة)", app_name)
    formula = ("Formula", "المعادلة",
               "Delivery fee = max(Minimum, Base + Per-km × Distance)",
               "سعر التوصيل = أكبر(الحد الأدنى, الأساسي + سعر/كم × المسافة)",
               COL["gold"])
    clarity = [
        ("Shown upfront", "تظهر قبل التأكيد", "Customer sees total fee before ordering", "العميل يرى الرسوم قبل تأكيد الطلب", COL["ok"]),
        ("Cash on delivery", "الدفع كاش", "Customer pays at delivery", "العميل يدفع عند الاستلام", COL["warn"]),
        ("Configurable", "قابل للضبط", "Rates can vary by zone/time", "الأسعار يمكن ضبطها حسب المنطقة/الوقت", COL["gold"]),
        ("Tips", "التيب", "Tips enabled for drivers", "التيب متاح للسائق", COL["gold"]),
    ]
    add_bilingual_cards(s, [formula] + clarity, Inches(0.9), Inches(1.55), Inches(12.0), Inches(5.4), cols=2)
    add_footer(s, idx, total); idx += 1

    # 11 Dispatch
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_header(s, "Dispatch & reliability", "الإسناد والموثوقية", app_name)
    chip(s, Inches(0.9), Inches(1.15), "Hybrid: auto + override", FONT_EN, COL["bg"], COL["gold"], w=Inches(3.2))
    chip(s, Inches(4.25), Inches(1.15), "No driver → Wait longer", FONT_EN, COL["bg"], COL["gold"], w=Inches(3.2))
    chip(s, Inches(7.6), Inches(1.15), "Live status updates", FONT_EN, COL["bg"], COL["gold"], w=Inches(2.7))
    add_dispatch_visual(s)
    bottom = [
        ("Dispatcher override", "تدخل الموزّع", "Manual assignment when needed", "تعيين يدوي عند الحاجة", COL["gold"]),
        ("Customer choice", "خيار العميل", "Wait longer instead of auto-cancel", "انتظار أكثر بدل الإلغاء التلقائي", COL["gold"]),
    ]
    add_bilingual_cards(s, bottom, Inches(0.9), Inches(4.9), Inches(12.0), Inches(1.8), cols=2)
    add_footer(s, idx, total); idx += 1

    # 12 Trust + Ops (split, no overlap)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_header(s, "Trust + operations", "الثقة + التشغيل", app_name)
    trust = [
        ("Driver verification", "تحقق السائق", "Phone + ID + selfie", "رقم الهاتف + هوية + سيلفي", COL["gold"]),
        ("Vehicle info", "بيانات المركبة", "Plate + vehicle type", "اللوحة + نوع المركبة", COL["gold"]),
        ("Proof of payment", "إثبات الدفع", "Receipt photo + merchant confirm", "صورة إيصال + تأكيد المتجر", COL["gold"]),
        ("Support", "الدعم", "In-app chat only", "محادثة داخل التطبيق فقط", COL["gold"]),
    ]
    add_bilingual_cards(s, trust, Inches(0.9), Inches(1.55), Inches(5.9), Inches(5.85), cols=1)
    add_admin_dashboard(s)
    add_footer(s, idx, total); idx += 1

    return prs

if __name__ == "__main__":
    prs = build_fixed_deck(app_name="Kaziony", city_en="Tripoli", city_ar="طرابلس")
    prs.save("Kaziony_Investor_Deck_FIXED.pptx")
